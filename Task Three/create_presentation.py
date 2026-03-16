"""Create OrderBot Presentation for Task 3."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 120, 212)
RED = RGBColor(220, 53, 69)
GREEN = RGBColor(40, 167, 69)
ORANGE = RGBColor(255, 153, 0)
LIGHT_GRAY = RGBColor(240, 240, 240)
GRAY = RGBColor(108, 117, 125)
DARK_BG = RGBColor(15, 23, 42)
CARD_BG = RGBColor(30, 41, 59)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, text="", font_size=11, font_color=WHITE, bold=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=13, font_color=WHITE, bold_first=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = (bold_first and i == 0)
        p.space_after = Pt(4)
    return txBox

def add_arrow(slide, left, top, width, height, color=GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_right_arrow(slide, left, top, width, height, color=GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ══════════════════════════════════════════════════════════════
# SLIDE 1: THE PROBLEM
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DARK_BG)

# Title bar
add_shape(slide1, 0, 0, 13.333, 1.1, DARK_BLUE)
add_text(slide1, 0.5, 0.15, 12, 0.8, "THE PROBLEM: Our Current Order Process Is a Bottleneck", font_size=28, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text(slide1, 0.5, 1.2, 12, 0.4, "Current Manual Hardware Order Processing Workflow", font_size=16, font_color=LIGHT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

# ── FLOWCHART ──
# Step 1: Email
add_shape(slide1, 0.8, 1.9, 2.2, 0.7, RGBColor(37, 99, 235), "1. Sales Rep Emails\nPDF Order Form", 11, WHITE, True)
add_arrow(slide1, 1.7, 2.6, 0.4, 0.35, RGBColor(100, 116, 139))

# Step 2: Manual Read
add_shape(slide1, 0.8, 3.05, 2.2, 0.7, RGBColor(37, 99, 235), "2. Ops Staff Manually\nReads PDF", 11, WHITE, True)
add_arrow(slide1, 1.7, 3.75, 0.4, 0.35, RGBColor(100, 116, 139))

# Step 3: Salesforce
add_shape(slide1, 0.8, 4.2, 2.2, 0.7, RGBColor(37, 99, 235), "3. Check Customer in\nSalesforce (Manual)", 11, WHITE, True)
add_arrow(slide1, 1.7, 4.9, 0.4, 0.35, RGBColor(100, 116, 139))

# Step 4: Google Sheets
add_shape(slide1, 0.8, 5.35, 2.2, 0.7, RGBColor(37, 99, 235), "4. Check Inventory in\nGoogle Sheets (Manual)", 11, WHITE, True)
add_arrow(slide1, 1.7, 6.05, 0.4, 0.35, RGBColor(100, 116, 139))

# Decision diamond - use a regular shape
add_shape(slide1, 0.8, 6.5, 2.2, 0.7, ORANGE, "5. DECISION:\nAll OK?", 11, WHITE, True, MSO_SHAPE.DIAMOND)

# OK path
add_right_arrow(slide1, 3.1, 5.5, 0.5, 0.3, GREEN)
add_shape(slide1, 3.7, 5.2, 2.0, 0.9, GREEN, "OK: Send 2 Emails\n- Customer Confirm\n- Warehouse Order", 10, WHITE, True)

# Problem path
add_right_arrow(slide1, 3.1, 6.65, 0.5, 0.3, RED)
add_shape(slide1, 3.7, 6.35, 2.0, 0.7, RED, "PROBLEM: Email\nSales Rep (Delays!)", 10, WHITE, True)

# ── PAIN POINTS (right side) ──
add_text(slide1, 6.5, 1.9, 6, 0.4, "Key Pain Points", font_size=20, font_color=RED, bold=True)

pain_points = [
    ("SLOW", "Each order takes 25-45 minutes to process manually"),
    ("ERROR-PRONE", "Manual data entry leads to ~8% error rate"),
    ("BOTTLENECK", "Single shared inbox — one person at a time"),
    ("NO AFTER-HOURS", "Orders received at night wait until next business day"),
    ("DELAYS", "Problems require back-and-forth emails (avg 4+ hours)"),
    ("UNSCALABLE", "Cannot handle volume spikes without more staff"),
]

y = 2.5
for title, desc in pain_points:
    add_shape(slide1, 6.5, y, 1.3, 0.55, RED, title, 10, WHITE, True)
    add_text(slide1, 7.9, y + 0.05, 5, 0.5, desc, font_size=12, font_color=RGBColor(203, 213, 225))
    y += 0.7

# Key Metrics bar at bottom
add_shape(slide1, 0.3, 7.05, 12.7, 0.35, RGBColor(30, 41, 59))
add_text(slide1, 0.5, 7.05, 12, 0.35, "Current Avg: 35 min/order  |  ~8% Error Rate  |  Business Hours Only (8am-5pm)  |  Max ~30 orders/day", font_size=12, font_color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: THE SOLUTION
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, DARK_BG)

# Title bar
add_shape(slide2, 0, 0, 13.333, 1.1, GREEN)
add_text(slide2, 0.5, 0.15, 12, 0.8, "THE SOLUTION: Introducing 'OrderBot' — Our Autonomous AI Agent", font_size=28, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide2, 0.5, 1.2, 12, 0.4, "Automated End-to-End Order Processing in Under 5 Minutes", font_size=16, font_color=GREEN, bold=False, alignment=PP_ALIGN.CENTER)

# ── AUTOMATED FLOWCHART (left side) ──
add_text(slide2, 0.5, 1.7, 5, 0.4, "How OrderBot Works", font_size=18, font_color=WHITE, bold=True)

steps = [
    ("1. DETECT", "Monitors inbox 24/7\nfor new order emails", RGBColor(16, 185, 129)),
    ("2. EXTRACT", "AI parses PDF and\nextracts order data", RGBColor(16, 185, 129)),
    ("3. VERIFY", "Auto-checks Salesforce\naccount status via API", RGBColor(16, 185, 129)),
    ("4. CHECK", "Auto-checks inventory\nin Google Sheets API", RGBColor(16, 185, 129)),
    ("5. DECIDE", "AI analyzes results:\nOK or Exception?", ORANGE),
]

y = 2.2
for title, desc, color in steps:
    add_shape(slide2, 0.5, y, 1.4, 0.65, color, title, 10, WHITE, True)
    add_text(slide2, 2.0, y + 0.05, 2.5, 0.6, desc, font_size=10, font_color=RGBColor(203, 213, 225))
    if y < 5.4:
        add_arrow(slide2, 1.05, y + 0.65, 0.3, 0.2, RGBColor(100, 116, 139))
    y += 0.85

# OK action
add_right_arrow(slide2, 2.0, 5.75, 0.4, 0.25, GREEN)
add_shape(slide2, 2.5, 5.55, 2.2, 0.65, GREEN, "AUTO-SEND\nConfirmation + Warehouse\n+ Update Salesforce", 9, WHITE, True)

# Exception action
add_right_arrow(slide2, 2.0, 6.55, 0.4, 0.25, ORANGE)
add_shape(slide2, 2.5, 6.35, 2.2, 0.55, ORANGE, "ESCALATE\nAlert ops team with\nfull context summary", 9, WHITE, True)

# ── AGENT ANATOMY (right side) ──
add_text(slide2, 5.5, 1.7, 7, 0.4, "OrderBot Agent Design", font_size=18, font_color=WHITE, bold=True)

components = [
    ("GOAL", "Process 100% of hardware orders from email to fulfillment\nconfirmation in under 5 minutes, with 99.5% accuracy.", GREEN),
    ("PERCEPTION\n(Input Tools)", "Email Inbox API — monitors for new orders 24/7\nPDF Parser — extracts customer, product, quantity data\nSalesforce API — reads customer account status\nGoogle Sheets API — reads real-time inventory levels", LIGHT_BLUE),
    ("PLANNING", "1. Detect new email  2. Parse PDF attachment\n3. Validate customer in Salesforce  4. Check inventory\n5. Decide: All OK → fulfill  |  Issue → escalate", ORANGE),
    ("ACTION\n(Output Tools)", "Email Sender — customer confirmation + warehouse order\nSalesforce API — updates order history automatically\nGoogle Sheets API — decrements inventory count", RGBColor(168, 85, 247)),
    ("MEMORY", "Logs every order for audit trail and analytics.\nLearns patterns: flags repeat-offender incomplete forms.\nBuilds knowledge base for continuous improvement.", RGBColor(236, 72, 153)),
]

y = 2.2
for title, desc, color in components:
    add_shape(slide2, 5.5, y, 1.8, 0.8, color, title, 9, WHITE, True)
    add_text(slide2, 7.4, y + 0.02, 5.5, 0.8, desc, font_size=9, font_color=RGBColor(203, 213, 225))
    y += 0.88


# ══════════════════════════════════════════════════════════════
# SLIDE 3: BUSINESS IMPACT
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, DARK_BG)

# Title bar
add_shape(slide3, 0, 0, 13.333, 1.1, DARK_BLUE)
add_text(slide3, 0.5, 0.15, 12, 0.8, "BUSINESS IMPACT & Recommended Next Steps", font_size=28, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ── METRICS ROW ──
add_text(slide3, 0.5, 1.2, 12, 0.4, "Projected Impact: Before vs After OrderBot", font_size=18, font_color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

metrics = [
    ("Processing Time", "35 min", "< 5 min", "85% Faster"),
    ("Error Rate", "~8%", "< 0.5%", "94% Reduction"),
    ("Availability", "8am–5pm", "24/7/365", "Always On"),
    ("Daily Capacity", "~30 orders", "500+ orders", "16x More"),
    ("Staff Hours/Week", "40+ hrs", "2 hrs (oversight)", "95% Saved"),
]

x = 0.4
for title, before, after, improvement in metrics:
    add_shape(slide3, x, 1.75, 2.4, 1.8, CARD_BG, "", 10, WHITE)
    add_text(slide3, x + 0.1, 1.8, 2.2, 0.3, title, font_size=12, font_color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide3, x + 0.1, 2.15, 2.2, 0.25, f"Before: {before}", font_size=10, font_color=RED, alignment=PP_ALIGN.CENTER)
    add_text(slide3, x + 0.1, 2.4, 2.2, 0.25, f"After: {after}", font_size=10, font_color=GREEN, alignment=PP_ALIGN.CENTER)
    add_text(slide3, x + 0.1, 2.7, 2.2, 0.6, improvement, font_size=18, font_color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    x += 2.55

# ── BUSINESS BENEFITS ──
add_text(slide3, 0.5, 3.8, 6, 0.4, "Key Business Benefits", font_size=18, font_color=GREEN, bold=True)

benefits = [
    "Cost Savings — Reduce manual processing costs by ~95%, freeing staff for higher-value work",
    "Speed — Orders processed in minutes, not hours. Customers get instant confirmation",
    "Accuracy — AI eliminates manual data entry errors, achieving 99.5% accuracy target",
    "Scalability — Handle volume spikes (seasonal, promotions) without hiring additional staff",
    "24/7 Operations — Process orders received at night, weekends, and holidays automatically",
    "Audit Trail — Every action logged for compliance, reporting, and continuous improvement",
]

y = 4.25
for b in benefits:
    add_text(slide3, 0.7, y, 6, 0.3, f"  {b}", font_size=11, font_color=RGBColor(203, 213, 225))
    y += 0.35

# ── NEXT STEPS (right side) ──
add_text(slide3, 7.2, 3.8, 5.5, 0.4, "Recommended Next Steps", font_size=18, font_color=ORANGE, bold=True)

add_shape(slide3, 7.2, 4.35, 5.6, 1.6, CARD_BG)

steps_text = [
    ("Phase 1: Read-Only Proof of Concept (Week 1-2)", GREEN),
    ("  OrderBot monitors inbox and processes orders in shadow mode", None),
    ("  Compares its decisions to human decisions — no actions taken", None),
    ("  Goal: Validate 99.5% accuracy before going live", None),
    ("", None),
    ("Phase 2: Supervised Rollout (Week 3-4)", ORANGE),
    ("  OrderBot processes orders with human approval for each action", None),
    ("  Staff reviews and approves/rejects before emails are sent", None),
    ("", None),
    ("Phase 3: Full Autonomous Mode (Week 5+)", LIGHT_BLUE),
    ("  OrderBot operates independently with exception escalation", None),
    ("  Weekly performance reviews and continuous improvement", None),
]

txBox = slide3.shapes.add_textbox(Inches(7.4), Inches(4.4), Inches(5.3), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, (text, color) in enumerate(steps_text):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = color if color else RGBColor(203, 213, 225)
    p.font.bold = True if color else False
    p.space_after = Pt(1)

# Bottom bar
add_shape(slide3, 0.3, 6.95, 12.7, 0.4, RGBColor(30, 41, 59))
add_text(slide3, 0.5, 6.95, 12, 0.4, "Recommendation: Begin Phase 1 (Read-Only PoC) immediately — zero risk, high insight, fast validation", font_size=13, font_color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# Save
prs.save("OrderBot_Presentation.pptx")
print("Presentation created: OrderBot_Presentation.pptx")
